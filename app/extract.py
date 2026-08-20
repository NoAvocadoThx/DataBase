"""资料文件文本提取 + 专家信息抽取（LLM 或规则）。"""
import json, os, re

import httpx

LLM_URL = os.getenv("LLM_URL", "https://api.deepseek.com/chat/completions")
LLM_KEY = os.getenv("LLM_API_KEY", "")
LLM_MODEL = os.getenv("LLM_MODEL", "deepseek-chat")

FIELDS = ["name", "org", "title", "field", "email", "phone", "topic", "source_text"]

PHONE_RE = re.compile(r"(?<!\d)1[3-9]\d{9}(?!\d)")
EMAIL_RE = re.compile(r"[\w.+-]+@[\w-]+(?:\.[\w-]+)+")
WECHAT_RE = re.compile(r"(?:微信|WeChat|wx)[:：\s]*([A-Za-z][\w-]{5,19})", re.I)
TITLE_WORDS = ["院士", "主任医师", "副主任医师", "主治医师", "教授", "副教授", "研究员", "副研究员",
               "主任", "副主任", "总经理", "副总裁", "总裁", "董事长", "首席科学家", "CEO", "CTO",
               "CMO", "博士", "审评员", "高级工程师"]
ORG_SUFFIX = r"(?:大学|学院|医院|研究所|研究院|公司|集团|中心|药监局|药审中心|协会|学会|基金会|科技|制药|医药|生物)"
CN_NAME = r"[一-龥]{2,4}"
ROLE_WORDS = {"主持人", "报告人", "致辞人", "嘉宾", "讲者", "演讲人", "主席", "开幕式", "讨论嘉宾", "点评人", "主持"}


def llm_enabled() -> bool:
    return bool(LLM_KEY)


# ---------- 文件 → 文本 ----------
def file_to_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        import fitz
        with fitz.open(path) as doc:
            return "\n".join(p.get_text() for p in doc)
    if ext == ".docx":
        import docx
        d = docx.Document(path)
        parts = [p.text for p in d.paragraphs]
        for t in d.tables:
            for row in t.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    if ext == ".pptx":
        from pptx import Presentation
        out = []
        for slide in Presentation(path).slides:
            for sh in slide.shapes:
                if sh.has_text_frame:
                    out.append(sh.text_frame.text)
        return "\n".join(out)
    with open(path, encoding="utf-8", errors="ignore") as f:
        return f.read()


# ---------- 脱敏 ----------
def redact(text: str) -> str:
    text = PHONE_RE.sub("[手机]", text)
    text = EMAIL_RE.sub("[邮箱]", text)
    return WECHAT_RE.sub("微信:[微信]", text)


# ---------- 规则抽取 ----------
def _find_title(seg: str) -> str:
    for w in TITLE_WORDS:
        if w in seg:
            return w
    return ""


def rule_extract(text: str) -> list[dict]:
    """逐行/逐句找 “姓名 + 单位(+职称)” 的模式。够用于议程、嘉宾名单等半结构化资料。"""
    results, seen = [], set()
    pat = re.compile(rf"({CN_NAME})[\s,，、:：|（(-]*([一-龥A-Za-z]*?{ORG_SUFFIX}[一-龥A-Za-z]*)")
    for line in re.split(r"[\n\r]+", text):
        line = line.strip()
        if not line or len(line) > 300:
            continue
        for m in pat.finditer(line):
            name, org = m.group(1), m.group(2)
            for rw in ROLE_WORDS:  # "主持人：王强" → 去掉角色前缀
                if name.startswith(rw) and len(name) > len(rw):
                    name = name[len(rw):]
            if name in seen or name in ROLE_WORDS or _find_title(name) or name.endswith(("大学", "医院", "公司")):
                continue
            seen.add(name)
            d = {k: "" for k in FIELDS}
            d.update(name=name, org=org, title=_find_title(line), source_text=line[:200])
            # 上下文内的联系方式
            em, ph = EMAIL_RE.search(line), PHONE_RE.search(line)
            if em:
                d["email"] = em.group()
            if ph:
                d["phone"] = ph.group()
            results.append(d)
    return results


# ---------- LLM 抽取 ----------
PROMPT = """从下面的会议/专家资料中提取所有出现的专家。只输出 JSON 数组，不要其他内容。
每个元素: {"name": 姓名, "org": 单位, "title": 职务或职称, "field": 研究方向, "topic": 报告主题(若有), "source_text": 原文中支持该条的一句话}
资料中无法确认的字段留空字符串，不要推断。

资料:
"""


def llm_extract(text: str) -> list[dict]:
    r = httpx.post(LLM_URL, timeout=90, headers={"Authorization": f"Bearer {LLM_KEY}"},
                   json={"model": LLM_MODEL, "temperature": 0,
                         "messages": [{"role": "user", "content": PROMPT + redact(text[:12000])}]})
    r.raise_for_status()
    txt = r.json()["choices"][0]["message"]["content"]
    arr = json.loads(re.search(r"\[.*\]", txt, re.S).group())
    out = []
    for it in arr:
        d = {k: "" for k in FIELDS}
        d.update({k: str(it.get(k, "") or "") for k in FIELDS})
        out.append(d)
    return out


def attach_local_contacts(text: str, cands: list[dict]) -> list[dict]:
    """手机/邮箱不经模型，在本地按姓名就近匹配（同一行）。"""
    lines = text.splitlines()
    for c in cands:
        if c.get("phone") and c.get("email"):
            continue
        for ln in lines:
            if c["name"] and c["name"] in ln:
                if not c.get("phone") and (m := PHONE_RE.search(ln)):
                    c["phone"] = m.group()
                if not c.get("email") and (m := EMAIL_RE.search(ln)):
                    c["email"] = m.group()
    return cands


def extract_experts(text: str) -> tuple[list[dict], str]:
    """返回 (候选列表, 使用的方式)。LLM 失败时回退规则。"""
    if llm_enabled():
        try:
            return attach_local_contacts(text, llm_extract(text)), "llm"
        except Exception as ex:  # 网络/解析失败不阻塞
            cands = rule_extract(text)
            return cands, f"rule(llm失败: {type(ex).__name__})"
    return rule_extract(text), "rule"
